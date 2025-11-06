from pptx import Presentation
from pptx.util import Inches, Pt
import os

PRJ_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
OUT_PATH = os.path.join(PRJ_ROOT, 'BehaviorLogger_Presentation.pptx')
FIG_DIR = os.path.join(PRJ_ROOT, 'docs', 'figures')

images = {
    'ui_main': os.path.join(FIG_DIR, 'ui_main.png'),
    'logger_terminal': os.path.join(FIG_DIR, 'logger_terminal.png'),
    'db_behavior_logs': os.path.join(FIG_DIR, 'db_behavior_logs.png'),
    'sql_schema': os.path.join(FIG_DIR, 'sql_schema.png'),
    'alert_and_csv': os.path.join(FIG_DIR, 'alert_and_csv.png'),
    'hasher_benchmark': os.path.join(FIG_DIR, 'hasher_benchmark.png'),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# helper to add notes
def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = notes_text

# helper to add title+bullets
def add_bullets_slide(title, bullets, notes=None, img_path=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(8.0), Inches(1.2), width=Inches(4.0))
        except Exception:
            pass
    if notes:
        add_notes(slide, notes)
    return slide

# Slide 1 - Title
layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
slide.shapes.title.text = "Behavior Logger — Real-time Telemetry & Anomaly Detection"
subtitle = slide.placeholders[1]
subtitle.text = "Single-host telemetry pipeline, Rust hasher, JDBC persistence, Java UI\nTeam: [Student Name(s)]  |  Date: November 2025"
add_notes(slide, "Introduce the project: goal, main components (Python collector, Rust hasher, MySQL, Java UI). Mention repo path: c:\\Users\\Ananya\\behavior_logger")

# Slide 2 - Motivation & Objectives
add_bullets_slide(
    "Motivation & Objectives",
    [
        "Motivation:\n- Near-real-time endpoint visibility, privacy, lightweight agent",
        "Objectives:\n- Capture keystrokes/mouse/app/metrics; persist hashed events",
        "- UI ~1 Hz refresh, processor every 5s, one-click CSV export"
    ],
    notes="Explain privacy (hashing) and resilience (JDBC primary, connector fallback). Mention SDG alignment."
)

# Slide 3 - Architecture
add_bullets_slide(
    "System Architecture (single-host)",
    [
        "Python Collector -> Rust Hasher -> MySQL <- Java Swing UI (JDBC)",
        "UI runs an internal processor every 5s to normalize and alert",
        "Fallbacks: Python SHA-256 and mysql-connector if needed"
    ],
    notes="Walk through data flow: capture -> hash -> persist -> UI/processor. Highlight fallback paths.",
    img_path=images.get('ui_main') if os.path.exists(images.get('ui_main', '')) else None
)

# Slide 4 - Collector: capture & hashing
add_bullets_slide(
    "Python Collector — Capture & Hashing",
    [
        "Input hooks: pynput (keyboard/mouse), Win32 for active window, psutil for system metrics",
        "Hashing: rust_hasher.exe (SHA-256); fallback to hashlib.sha256()",
        "Persistence: JPype/jaydebeapi (JDBC) first, then mysql-connector-python"
    ],
    notes="Show the call to rust_hasher and fallback path. Mention argv vs stdin tradeoff.",
    img_path=images.get('logger_terminal') if os.path.exists(images.get('logger_terminal', '')) else None
)

# Slide 5 - Storage & Schema
add_bullets_slide(
    "Database Schema & Persistence",
    [
        "Main table: behavior_logs (id, timestamp, event_type, hashed_event, raw_data, prediction, processed_at)",
        "Normalized tables: keystroke_events, mouse_events, application_events, system_metrics, alerts, anomaly_stats",
        "Resilience: JDBC primary, mysql-connector fallback, SKIP_DB_MODE when DB unavailable"
    ],
    notes="Explain that behavior_logs is a durable sink and normalization is performed by the UI processor.",
    img_path=images.get('sql_schema') if os.path.exists(images.get('sql_schema', '')) else None
)

# Slide 6 - UI & Real-time Processing
add_bullets_slide(
    "Java Swing UI — Live View & Processor",
    [
        "Features: table view, analytics, live log, JFreeChart time-series, CSV export",
        "Real-time: 1s auto-refresh; incremental using lastSeenId/lastSeenTimestamp",
        "Processor: processLogsInUI() runs every 5s and normalizes rows into domain tables"
    ],
    notes="Explain trimming of UI buffers (2000 rows, chart cap 600) and how start/stop logging works.",
    img_path=images.get('ui_main') if os.path.exists(images.get('ui_main', '')) else None
)

# Slide 7 - Classification Logic
add_bullets_slide(
    "Anomaly Classification — Heuristics & Alerts",
    [
        "Heuristics run in processLogsInUI(): examples: SYSTEM in type -> anomaly; raw_data containing 'CPU' -> anomaly",
        "Actions: set prediction (0 normal, 1 anomaly), insert alerts, insert behavior_events, upsert anomaly_stats",
        "Limitations: heuristic, not ML; next steps: local lightweight model or realtime_infer.py"
    ],
    notes="Describe anomaly pipeline and what prediction values mean. Suggest model upgrade path.",
    img_path=images.get('alert_and_csv') if os.path.exists(images.get('alert_and_csv', '')) else None
)

# Slide 8 - Performance & Trade-offs
add_bullets_slide(
    "Performance, Safety & Trade-offs",
    [
        "Rust hasher provides native speed; subprocess vs pyo3 tradeoffs for latency",
        "Memory & safety: Python GC + Rust ownership; bounded UI buffers and DB LIMITs",
        "Deployment: single-host prototype; credentials need externalization"
    ],
    notes="Explain tradeoffs and recommended mitigations: prebuilt binaries, pyo3 extension or persistent worker, remove hard-coded creds."
)

# Slide 9 - Demo & Evidence
add_bullets_slide(
    "Demo & Evidence",
    [
        "Run: setup_database(), run python/behavior_log.py, launch Java UI, Start Logging",
        "Observed: events in behavior_logs, live UI updates, alerts generated by processor",
        "Screenshots: logger_terminal.png, db_behavior_logs.png, ui_main.png"
    ],
    notes="Walk through demo steps and indicate where screenshots are located in docs/figures.",
    img_path=images.get('db_behavior_logs') if os.path.exists(images.get('db_behavior_logs', '')) else None
)

# Slide 10 - Next Steps & Conclusion
add_bullets_slide(
    "Conclusions & Next Steps",
    [
        "Achieved: end-to-end pipeline with hashing, persistence, UI and alerts",
        "Shortcomings: heuristic detection, single-host scope, subprocess overhead",
        "Next steps: pyo3 or persistent worker, lightweight ML inference, harden deployment, benchmark"
    ],
    notes="Summarize accomplishments and prioritized next steps. Invite reviewers to try the repo and appendix.",
)

# Save presentation
prs.save(OUT_PATH)
print(f"Saved PPTX to: {OUT_PATH}")
